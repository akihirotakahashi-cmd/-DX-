"""
添付ファイルからテキストを抽出するサービス。
対応形式: PDF / Word(.docx) / Excel(.xlsx) / PowerPoint(.pptx) /
          テキスト系(.txt, .csv, .md, .json, .xml)
"""
import io
from pathlib import Path

_MAX_CHARS = 15000


def extract_text_from_file(filename: str, content: bytes) -> str:
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(content)[:_MAX_CHARS]
        elif ext == ".docx":
            return _extract_docx(content)[:_MAX_CHARS]
        elif ext in (".xlsx", ".xls"):
            return _extract_xlsx(content)[:_MAX_CHARS]
        elif ext in (".pptx", ".ppt"):
            return _extract_pptx(content)[:_MAX_CHARS]
        elif ext in (".txt", ".md", ".csv", ".tsv", ".json", ".xml", ".html", ".htm"):
            return content.decode("utf-8", errors="replace")[:_MAX_CHARS]
        else:
            # バイナリ以外はとりあえずテキストとして試みる
            try:
                return content.decode("utf-8")[:_MAX_CHARS]
            except Exception:
                return f"[未対応フォーマット: {filename}]"
    except Exception as e:
        return f"[読み込みエラー ({filename}): {e}]"


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text)
    return "\n\n".join(pages)


def _extract_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    lines: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            lines.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.replace("|", "").strip():
                lines.append(row_text)
    return "\n".join(lines)


def _extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"=== {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join("" if v is None else str(v) for v in row)
            if row_text.replace("|", "").strip():
                lines.append(row_text)
    return "\n".join(lines)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"=== スライド {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)
